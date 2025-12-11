import fitz  
from pptx import Presentation
from docx import Document
from pdf2image import convert_from_path
import easyocr
import os


ocr_reader = easyocr.Reader(['en'], gpu=False)



def is_scanned_pdf(pdf_path):
    with fitz.open(pdf_path) as doc:
        for page in doc:
            if page.get_text():
                return False  
    return True  



def extract_pdf_text(pdf_path):
    text = ""
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text += page.get_text()
    return text.strip()



def extract_scanned_pdf(pdf_path):
    images = convert_from_path(pdf_path)
    all_text = ""

    for i, img in enumerate(images):
        img_path = f"temp_page_{i}.png"
        
        
        img.thumbnail((1024, 1024))
        img.save(img_path)

        
        page_text = extract_image(img_path)
        all_text += f"\n--- Page {i+1} ---\n{page_text}\n"

        if os.path.exists(img_path):
            os.remove(img_path)

    return all_text.strip()



def extract_image(image_path):
    
    ocr_result = ocr_reader.readtext(image_path, detail=0)
    ocr_text = "\n".join(ocr_result).strip()
    
    
    try:
        from llm_client import call_vision
        print(f"DEBUG: Analyze image {image_path}")
        vision_desc = call_vision("Describe this educational image in detail. Mention title, axes, data points, and relationships.", image_path)
        
        final_text = f"{ocr_text}\n\n[Visual Description]:\n{vision_desc}"
        return final_text.strip()
    except Exception as e:
        print(f"Vision extraction failed: {e}")
        return ocr_text



def extract_docx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([p.text for p in doc.paragraphs])
    return text.strip()



def extract_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text.strip()



def extract_txt(txt_path):
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()



def extract_text(file_path):
    extension = file_path.lower().split(".")[-1]

    if extension == "pdf":
        if is_scanned_pdf(file_path):
            return extract_scanned_pdf(file_path)
        return extract_pdf_text(file_path)

    elif extension == "docx":
        return extract_docx(file_path)

    elif extension == "pptx":
        return extract_pptx(file_path)

    elif extension in ["png", "jpg", "jpeg"]:
        return extract_image(file_path)

    elif extension == "txt":
        return extract_txt(file_path)

    else:
        return "Unsupported file format"
